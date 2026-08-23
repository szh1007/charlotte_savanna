"""输入解析器测试 (Issue 07 阶段 1: parsing).

覆盖: 文本归一化 / 标题提取 / 多格式文件解析 (txt/md/html/pdf/docx/pptx,
内存构造样本) / 链接抓取 (mock httpx) / 不支持格式抛错.
"""

import io

import pytest

from project.charplot.pipeline import parsers


def test_normalize_text_collapses_blank_lines():
    text = " 第一行  \n\n\n  第二行  \n\n第三行"
    assert parsers.normalize_text(text) == "第一行\n\n第二行\n\n第三行"


def test_extract_title_uses_first_line():
    text = "Python 装饰器指南\n\n第一章"
    assert parsers.extract_title(text) == "Python 装饰器指南"
    assert parsers.extract_title("  \n  ", fallback="默认") == "默认"


def test_parse_txt_utf8_and_gbk():
    assert parsers.parse_document("a.txt", "你好世界".encode()) == "你好世界"
    assert parsers.parse_document("a.txt", "中文内容".encode("gbk")) == "中文内容"


def test_parse_markdown():
    data = "# 标题\n\n正文段落".encode()
    # markdown 解析保留标题语法 (#), 仅归一化空白
    assert parsers.parse_document("a.md", data) == "# 标题\n\n正文段落"


def test_parse_html_strips_scripts():
    data = (
        "<html><head><script>var x=1;</script></head><body>"
        "<h1>标题</h1><nav>导航</nav><p>正文</p></body></html>"
    ).encode()
    text = parsers.parse_document("a.html", data)
    assert "标题" in text and "正文" in text
    assert "script" not in text and "导航" not in text


def test_parse_docx():
    import docx

    doc = docx.Document()
    doc.add_paragraph("第一段内容")
    doc.add_paragraph("第二段内容")
    buf = io.BytesIO()
    doc.save(buf)
    text = parsers.parse_document("a.docx", buf.getvalue())
    assert "第一段内容" in text and "第二段内容" in text


def test_parse_pptx():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "幻灯片标题"
    slide.placeholders[1].text = "幻灯片正文"
    buf = io.BytesIO()
    prs.save(buf)
    text = parsers.parse_document("a.pptx", buf.getvalue())
    assert "幻灯片标题" in text and "幻灯片正文" in text


def test_parse_pdf():
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    # pypdf 写出的空白页无文本, 断言不抛错且产出空字符串即可 (解析链路通)
    text = parsers.parse_document("a.pdf", buf.getvalue())
    assert isinstance(text, str)


def test_parse_unsupported_extension():
    with pytest.raises(parsers.ParseError, match="不支持的文件格式"):
        parsers.parse_document("a.xyz", b"data")


def test_fetch_link_content(monkeypatch):
    class FakeResp:
        content = "<html><body><h1>标题</h1><p>正文</p></body></html>".encode()

        def raise_for_status(self):
            pass

    class FakeClient:
        @staticmethod
        def get(url, timeout=None, follow_redirects=False, headers=None):
            return FakeResp()

    monkeypatch.setattr("httpx.get", FakeClient.get)
    text = parsers.fetch_link_content("https://example.com")
    assert "正文" in text


def test_fetch_link_failure_raises_parse_error(monkeypatch):
    import httpx

    def boom(url, **kwargs):
        raise httpx.ConnectError("connection refused")

    monkeypatch.setattr("httpx.get", boom)
    with pytest.raises(parsers.ParseError, match="链接抓取失败"):
        parsers.fetch_link_content("https://example.com")
