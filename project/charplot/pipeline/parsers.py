"""输入归一化解析 (Issue 07 阶段 1: parsing).

统一管道 (ADR-0002) 的输入侧: 文本归一化 / 网页链接抓取 / 文件解析
(pdf/docx/pptx/md/txt/html, 复用项目文档解析库). 全部产出一份纯文本
材料, 供后续分析/解构使用.

file 输入的文件二进制经 Django 内部端点 (CONTRACT.md §5) 获取, 解析
在 FastAPI 侧完成 (AI 能力端职责).
"""

import io
import logging

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt": "text",
    ".md": "markdown",
    ".html": "html",
    ".htm": "html",
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
}


class ParseError(ValueError):
    """解析失败 (文件格式不支持 / 内容为空 / 链接抓取失败)."""


def normalize_text(content: str) -> str:
    """文本归一化: 去行首尾空白 + 压缩连续空行为单个, 控制解析量上限."""
    lines = [line.strip() for line in content.splitlines()]
    text = "\n".join(lines).strip()
    while "\n\n\n" in text:
        text = text.replace("\n\n\n", "\n\n")
    # 超长材料截断: 后文由 LLM 分析处理, 完整原文无必要全量保留
    return text[:200_000]


def fetch_link_content(url: str) -> str:
    """抓取网页正文 (链接输入). 网络/解析失败抛 ParseError (任务可重试)."""
    import httpx
    from bs4 import BeautifulSoup

    from ..api import config

    try:
        resp = httpx.get(
            url,
            timeout=config.LINK_FETCH_TIMEOUT,
            follow_redirects=True,
            headers={"User-Agent": "Mozilla/5.0 (CharPlot pipeline)"},
        )
        resp.raise_for_status()
    except httpx.HTTPError as exc:
        raise ParseError(f"链接抓取失败 ({url}): {exc}") from exc
    soup = BeautifulSoup(resp.content, "html.parser")
    # 移除脚本/样式/导航噪声, 保留正文文本
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav"]):
        tag.decompose()
    text = soup.get_text("\n", strip=True)
    text = normalize_text(text)
    if not text:
        raise ParseError(f"链接正文为空 ({url})")
    return text


def _decode_bytes(data: bytes) -> str:
    """文本类文件编码探测: 优先 utf-8, 失败回退 gbk (中文环境)."""
    for encoding in ("utf-8", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _parse_pdf(data: bytes) -> str:
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(data))
    pages = [(page.extract_text() or "") for page in reader.pages]
    return normalize_text("\n".join(pages))


def _parse_docx(data: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return normalize_text("\n".join(parts))


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            parts.append("\n".join(texts))
    return normalize_text("\n\n".join(parts))


def _parse_html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav"]):
        tag.decompose()
    return normalize_text(soup.get_text("\n", strip=True))


def parse_document(filename: str, data: bytes) -> str:
    """按扩展名分发解析文件二进制, 返回纯文本材料.

    支持 .txt/.md/.html/.pdf/.docx/.pptx (PRD B-1); 其他格式抛 ParseError.
    """
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    ext = f".{ext}"
    kind = SUPPORTED_EXTENSIONS.get(ext)
    if kind is None:
        raise ParseError(f"不支持的文件格式: {ext or '(无扩展名)'}")

    try:
        if kind in ("text", "markdown"):
            text = _decode_bytes(data)
            return normalize_text(text)
        if kind == "html":
            return _parse_html(data)
        if kind == "pdf":
            return _parse_pdf(data)
        if kind == "docx":
            return _parse_docx(data)
        if kind == "pptx":
            return _parse_pptx(data)
    except ParseError:
        raise
    except Exception as exc:
        logger.warning("文件解析异常 (%s): %s", filename, exc)
        raise ParseError(f"文件解析失败 ({filename}): {exc}") from exc

    raise ParseError(f"文件解析失败 ({filename})")


def extract_title(text: str, fallback: str = "学习材料") -> str:
    """从材料文本提取标题 (首行截断, 与 Django 侧 derive 规则一致)."""
    for line in text.splitlines():
        line = line.strip()
        if line:
            return line[:50]
    return fallback
